# import pyperclip
# from ruamel.yaml import YAML
from xml.etree.ElementTree import XMLParser, parse
from enum import IntFlag
from dataclasses import dataclass, field, asdict
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx
import re
from pyscript import display, document
from js import loadpptx
import js
import io
from datetime import datetime
from pptx import Presentation
from pprint import pprint
import json
FILEINPUT = document.getElementById("inputpptx")
YAMLOUTPUT = document.getElementById("output")


async def translate_pptx(event):
    # print("test", dir(event))
    # print(dir(FILEINPUT.files))
    if FILEINPUT.files.length < 1:
        print("No file loaded.")

    if FILEINPUT.files.length > 0:
        content = io.BytesIO((await loadpptx()).encode('latin1'))
        # (await FILEINPUT.files.item(0).text()).encode('utf8'))
        presentation = Presentation(content)
        YAMLOUTPUT.innerText = divomathyaml(presentation)


# yaml = YAML()

OXMLSCHEMAREGEX = re.compile(
    "{http://schemas.openxmlformats.org/(.*)/2006/main}")

OXMLTAGS = {"avLst": "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst",  # {{{
            "blip": "{http://schemas.openxmlformats.org/drawingml/2006/main}blip",
            "blipFill": "{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill",
            "bodyPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr",
            "br": "{http://schemas.openxmlformats.org/drawingml/2006/main}br",
            "buChar": "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar",
            "buClrTx": "{http://schemas.openxmlformats.org/drawingml/2006/main}buClrTx",
            "buFont": "{http://schemas.openxmlformats.org/drawingml/2006/main}buFont",
            "buFontTx": "{http://schemas.openxmlformats.org/drawingml/2006/main}buFontTx",
            "buNone": "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone",
            "buSzTx": "{http://schemas.openxmlformats.org/drawingml/2006/main}buSzTx",
            "cNvGrpSpPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvGrpSpPr",
            "cNvPicPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPicPr",
            "cNvPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr",
            "cNvSpPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvSpPr",
            "chExt": "{http://schemas.openxmlformats.org/drawingml/2006/main}chExt",
            "chOff": "{http://schemas.openxmlformats.org/drawingml/2006/main}chOff",
            "creationId": "{http://schemas.microsoft.com/office/drawing/2014/main}creationId",
            "cs": "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
            "defRPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr",
            "ea": "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
            "effectLst": "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
            "effectRef": "{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef",
            "endParaRPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr",
            "ext": "{http://schemas.openxmlformats.org/drawingml/2006/main}ext",
            "extLst": "{http://schemas.openxmlformats.org/drawingml/2006/main}extLst",
            "fillRect": "{http://schemas.openxmlformats.org/drawingml/2006/main}fillRect",
            "fillRef": "{http://schemas.openxmlformats.org/drawingml/2006/main}fillRef",
            "fontRef": "{http://schemas.openxmlformats.org/drawingml/2006/main}fontRef",
            "gd": "{http://schemas.openxmlformats.org/drawingml/2006/main}gd",
            "grpSp": "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp",
            "grpSpPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr",
            "hiddenFill": "{http://schemas.microsoft.com/office/drawing/2010/main}hiddenFill",
            "latin": "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
            "ln": "{http://schemas.openxmlformats.org/drawingml/2006/main}ln",
            "lnRef": "{http://schemas.openxmlformats.org/drawingml/2006/main}lnRef",
            "lnSpc": "{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc",
            "lstStyle": "{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle",
            "lumMod": "{http://schemas.openxmlformats.org/drawingml/2006/main}lumMod",
            "lumOff": "{http://schemas.openxmlformats.org/drawingml/2006/main}lumOff",
            "lvl1pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr",
            "lvl2pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl2pPr",
            "lvl3pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl3pPr",
            "lvl4pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl4pPr",
            "lvl5pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl5pPr",
            "lvl6pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl6pPr",
            "lvl7pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl7pPr",
            "lvl8pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl8pPr",
            "lvl9pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}lvl9pPr",
            "noFill": "{http://schemas.openxmlformats.org/drawingml/2006/main}noFill",
            "normAutofit": "{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit",
            "nvGrpSpPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr",
            "nvPicPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr",
            "nvPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr",
            "nvSpPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr",
            "off": "{http://schemas.openxmlformats.org/drawingml/2006/main}off",
            "p": "{http://schemas.openxmlformats.org/drawingml/2006/main}p",
            "pPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr",
            "ph": "{http://schemas.openxmlformats.org/presentationml/2006/main}ph",
            "pic": "{http://schemas.openxmlformats.org/presentationml/2006/main}pic",
            "picLocks": "{http://schemas.openxmlformats.org/drawingml/2006/main}picLocks",
            "prstGeom": "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom",
            "r": "{http://schemas.openxmlformats.org/drawingml/2006/main}r",
            "rPr": "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr",
            "schemeClr": "{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr",
            "shade": "{http://schemas.openxmlformats.org/drawingml/2006/main}shade",
            "solidFill": "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill",
            "sp": "{http://schemas.openxmlformats.org/presentationml/2006/main}sp",
            "spAutoFit": "{http://schemas.openxmlformats.org/drawingml/2006/main}spAutoFit",
            "spLocks": "{http://schemas.openxmlformats.org/drawingml/2006/main}spLocks",
            "spPr": "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr",
            "spcAft": "{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft",
            "spcBef": "{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef",
            "spcPct": "{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct",
            "spcPts": "{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts",
            "srcRect": "{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect",
            "srgbClr": "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
            "stretch": "{http://schemas.openxmlformats.org/drawingml/2006/main}stretch",
            "style": "{http://schemas.openxmlformats.org/presentationml/2006/main}style",
            "sym": "{http://schemas.openxmlformats.org/drawingml/2006/main}sym",
            "t": "{http://schemas.openxmlformats.org/drawingml/2006/main}t",
            "tabLst": "{http://schemas.openxmlformats.org/drawingml/2006/main}tabLst",
            "txBody": "{http://schemas.openxmlformats.org/presentationml/2006/main}txBody",
            "useLocalDpi": "{http://schemas.microsoft.com/office/drawing/2010/main}useLocalDpi",
            "xfrm": "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
            }  # }}}


@dataclass
class BoundingBox:
    x: int = 0
    y: int = 0
    w: int = 0
    h: int = 0


@dataclass
class Attributes:
    empty: str = ""


@dataclass
class TextAttributes(Attributes):
    size: int = field(default=21)
    weight: str = "normal"
    color: str = "#000000"
    alignment: str = "left"


@dataclass
class FrameAttributes(Attributes):
    border_style: str = "solid"
    border_color: str = "#000"
    border_width: int = field(default=1)
    background_color: str = "#fff"
    border_radius: int = field(default=0)


@dataclass
class DivomathComponent:
    type: str
    name: str
    boundingBox: BoundingBox
    valueType: str
    uuid: str
    media: list[str] = field(default_factory=list)
    value: list[str] = field(default_factory=list)
    attributes: Attributes = field(default_factory=Attributes)
    state: list = field(default_factory=list)


def pptx2divomath(presentation):
    info = collect_presentation_info(presentation)
    slide_width = info["presentation_properties"]["slide_width"]
    slide_height = info["presentation_properties"]["slide_height"]
    all_components = []
    for elements in info["slides"]:
        components = []
        for i, element in enumerate(elements):
            if element["type"] == ElementType.TEXT:
                # @refactor: BoundingBox.from_element(...)
                bbox = BoundingBox(
                    round(element["shape"]["left"]/slide_width*100, 2),
                    round(element["shape"]["top"]/slide_height*100, 2),
                    round(element["shape"]["width"]/slide_width*100, 2),
                    round(element["shape"]["height"]/slide_height*100, 2))
                component_value = []
                # @improve: for now run = paragraph
                for paragraph in element["text"]["paragraphs"]:
                    print(paragraph)
                    for run in paragraph["runs"]:
                        component_value.append(
                            f'<p><span style=\"font-size: {run["font"]["size"]};\">{run["content"]}</span></p>')
                        components.append(DivomathComponent("TEXT", f"TEXT-{i}", bbox,
                                                            "RICHTEXT", "eid-15", value=component_value,
                                                            attributes=TextAttributes(alignment=run["alignment"])))
            elif element["type"] == ElementType.IMAGE:

                print("image")
                bbox = BoundingBox(
                    round(element["shape"]["left"]/slide_width*100, 2),
                    round(element["shape"]["top"]/slide_height*100, 2),
                    round(element["shape"]["width"]/slide_width*100, 2),
                    round(element["shape"]["height"]/slide_height*100, 2))
                components.append(DivomathComponent("IMAGE", f"IMAGE-{i}", bbox,
                                                    "UUID", "eid-15", value=["0"],
                                                    attributes=Attributes()))
            elif element["type"] == ElementType.SHAPE:
                print("shape")
                bbox = BoundingBox(
                    round(element["shape"]["left"]/slide_width*100, 2),
                    round(element["shape"]["top"]/slide_height*100, 2),
                    round(element["shape"]["width"]/slide_width*100, 2),
                    round(element["shape"]["height"]/slide_height*100, 2))
                components.append(DivomathComponent("FRAME", f"FRAME-{i}", bbox,
                                                    "NULL", "eid-15", value=["0"],
                                                    attributes=FrameAttributes(
                                                        border_style="solid",
                                                        border_color="#000",
                                                        border_width=1,
                                                        background_color="#fff",
                                                        border_radius=0
                                                    )))

        all_components.append(components)
    return all_components


def walk_xmlelement(element, depth=None):
    children = []
    next_depth = None
    if depth is not None:
        if depth == 0:
            return "..."
        next_depth = depth - 1
    for child in element:
        tag_subst = child.tag
        tag_match = OXMLSCHEMAREGEX.match(child.tag)
        if tag_match is not None:
            tag_subst = child.tag.replace(
                child.tag[slice(*tag_match.span())], "")
        text = None
        # only if it contains more than whitespace characters
        if child.text is not None and child.text.strip():
            text = child.text
        children.append({'tag': tag_subst, 'attrib': child.attrib, 'data': text,
                        'children': walk_xmlelement(child, next_depth)})
    return children


def attribute_overview(obj):
    attributes = {}
    for key in dir(obj):
        value = None
        try:
            value = obj.__getattribute__(key)
        except:
            pass
        attributes[key] = value, str(value)
    return attributes


class ElementType(IntFlag):

    SHAPE = 0
    TEXT = 1
    IMAGE = 2


@ dataclass
class SlideElements:
    type: None
    image: None
    details: dict


class Length:
    def __init__(self, emu):
        self.emu = emu
        self.inch = emu/pptx.util.Inches(1)
        self.cm = emu/pptx.util.Cm(1)
        self.mm = emu/pptx.util.Mm(1)
        self.pt = emu/pptx.util.Pt(1)
        self.cpt = emu/pptx.util.Centipoints(1)

    def __add__(self, other):
        if isinstance(other, Length):
            return Length(self.emu + other.emu)
        else:
            raise NotImplementedError(f"not implemented for {type(other)}")

    def __sub__(self, other):
        if isinstance(other, Length):
            return Length(self.emu - other.emu)
        else:
            raise NotImplementedError(f"not implemented for {type(other)}")

    def __mul__(self, other):
        if isinstance(other, int):
            return Length(self.emu * other)
        else:
            raise NotImplementedError(f"not implemented for {type(other)}")

    def __rmul__(self, other):
        if isinstance(other, int):
            return Length(self.emu * other)
        else:
            raise NotImplementedError(f"not implemented for {type(other)}")

    def __truediv__(self, other):
        if isinstance(other, Length):
            return self.emu / other.emu
        elif isinstance(other, int):
            return Length(self.emu // other)
        else:
            raise NotImplementedError(f"not implemented for {type(other)}")

    def __rtruediv__(self, other):
        if isinstance(other, Length):
            return other.emu / self.emu
        else:
            raise NotImplementedError(f"not implemented for {type(other)}")

    def __repr__(self):
        return f"Length({self.emu})"

    def __str__(self):
        return f"Length(emu={self.emu},mm={self.mm},inch={self.inch})"


def collect_presentation_info(presentation, slides_fromto=None):

    from_slide = 0
    to_slide = len(presentation.slides)

    if hasattr(slides_fromto, "__iter__"):
        from_slide, to_slide = slides_fromto[0] - 1, slides_fromto[1] - 1

    info = {
        "presentation_properties": {
            "slide_height": Length(presentation.slide_height),
            "slide_width": Length(presentation.slide_width),
            "num_slides_total": len(presentation.slides)
        },

        "num_slides": to_slide - from_slide + 1,
        "slides": [get_slide_elements(slide)
                   for i, slide in enumerate(presentation.slides)
                   if from_slide <= i <= to_slide]

    }

    return info


def get_slide_elements(slide):
    # 'adjustments',
    # 'auto_shape_type',
    # 'click_action',
    # 'fill',
    # 'get_or_add_ln',
    # 'has_chart',
    # 'has_table',
    # 'has_text_frame',
    # 'is_placeholder',
    # 'line',
    # 'ln',
    # 'name',
    # 'part',
    # 'placeholder_format',
    # 'shadow',
    # 'shape_id',
    # 'shape_type',
    # 'text_frame',
    slide_elements = []
    for shape in slide.shapes:
        element_type = ElementType.SHAPE
        geometry = shape.element.find(f'.//{OXMLTAGS["prstGeom"]}')
        if geometry is not None:
            geometry = geometry.attrib["prst"]
        shape_style = {
            "top": Length(shape.top),
            "left": Length(shape.left),
            "height": Length(shape.height),
            "width": Length(shape.width),
            "rotation": shape.rotation,
            "geometry": geometry,
            "shape_on_slide_id": shape.shape_id,
            "shape_type": str(shape.shape_type),
            "is_placeholder": shape.is_placeholder
        }

        text = None
        if shape.has_text_frame and shape.text:  # {{{
            element_type = ElementType.TEXT

            text = {"content": shape.text,
                    "margins": {"top": Length(shape.text_frame.margin_top),
                                "right": Length(shape.text_frame.margin_right),
                                "bottom": Length(shape.text_frame.margin_bottom),
                                "left": Length(shape.text_frame.margin_left)},
                    "is_wrapping": shape.text_frame.word_wrap,
                    "paragraphs": []
                    }
            for para in shape.text_frame.paragraphs:
                info = {}
                # @TODO: Better name?

                alignment = "left"
                if para.alignment in [1, 2, 3]:
                    alignment = ["", "left", "center", "right"][para.alignment]
                info["runs"] = []
                for run in para.runs:

                    if run.font.size is None:
                        # @improve: hardcoded value
                        fontsize = 22
                    else:
                        fontsize = int(Length(run.font.size).pt)
                    # import IPython
                    # IPython.embed()
                    info["runs"].append({
                        "content": run.text,
                        "alignment": alignment,
                        "font": {
                            "size": fontsize,
                            "name": run.font.name,
                            "is_bold": run.font.bold,
                            "is_italic": run.font.italic,
                            "is_underlined": run.font.underline,
                        },
                        "link": run.hyperlink.address
                    })

                text["paragraphs"].append(info)
# }}}

        # @improve: use Enum
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print("image", shape.has_text_frame)
            element_type = ElementType.IMAGE

        xml = shape.element
        xmldict = walk_xmlelement(xml)
        cNvPr = shape.element.find(f'.//{OXMLTAGS["cNvPr"]}')
        description = None
        if "descr" in cNvPr.attrib:
            description = cNvPr.attrib["descr"]
        # @fix: the js object-esque function => custom dataclass instead of list
        slide_elements.append({'type': element_type,
                               'name': shape.name,
                               'shape': shape_style,
                               'text': text,
                               'description': description,
                               'xml': lambda _: xmldict})
    return slide_elements


def show_slides(presentation):
    print(presentation.slides)


def example_pptx(key="all"):
    slides = {"one": "OneSlide",
              "all": "_2024_Malaufgaben am Punktefeld darstellen-1_240909"}
    return Presentation(f"../examples/{slides[key]}.pptx")


def divomathyaml(presentation):
    s = io.StringIO()
    s.write("[")
    for i, it in enumerate(pptx2divomath(presentation)[0]):
        if not i == 0:
            s.write(",")
        s.write(json.dumps(asdict(it)).replace(r'\"', r'\\\"'))
    s.write("]")
    # [yaml.dump([asdict(it)], s) for it in pptx2divomath(presentation)[0]]
    return s.getvalue()
